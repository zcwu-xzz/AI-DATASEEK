import importlib.util
from pathlib import Path

from pptx import Presentation

ROOT=Path(__file__).resolve().parents[2]
SPEC=importlib.util.spec_from_file_location("presentation_operations",ROOT/"tools"/"presentations"/"operations.py")
OPS=importlib.util.module_from_spec(SPEC); assert SPEC.loader is not None; SPEC.loader.exec_module(OPS)


def make_deck(path):
    deck=Presentation(); slide=deck.slides.add_slide(deck.slide_layouts[5]); slide.shapes.title.text="Climate summary"
    table=slide.shapes.add_table(2,2,0,1000000,4000000,1500000).table
    table.cell(0,0).text="region"; table.cell(0,1).text="value"; table.cell(1,0).text="A"; table.cell(1,1).text="1"
    deck.save(path)


def test_pptx_inspect_and_extract(tmp_path, monkeypatch):
    source=tmp_path/"sample.pptx"; make_deck(source); monkeypatch.setenv("AI_DATASEEK_OUTPUT_ROOT",str(tmp_path))
    inspected=OPS.inspect_presentation({"input_path":str(source)})
    assert inspected["summary"]["slide_count"] == 1
    assert inspected["summary"]["slides"][0]["table_count"] == 1
    extracted=OPS.extract_content({"input_path":str(source),"max_chars":10000,"output_path":str(tmp_path/"slides.md")})
    assert extracted["summary"]["slides"][0]["title"] == "Climate summary"
    assert extracted["artifacts"][0]["size_bytes"] > 0
