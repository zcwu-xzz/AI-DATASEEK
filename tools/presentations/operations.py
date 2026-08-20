#!/usr/bin/env python3
from __future__ import annotations

import argparse
import base64
import json
import math
import os
import subprocess
import tempfile
from pathlib import Path
from typing import Any

import numpy as np


def clean(value: Any) -> Any:
    if value is None or isinstance(value, (str, bool, int)): return value
    if isinstance(value, (float, np.floating)): return float(value) if math.isfinite(float(value)) else None
    if isinstance(value, np.integer): return int(value)
    if isinstance(value, dict): return {str(k): clean(v) for k, v in value.items()}
    if isinstance(value, (list, tuple, np.ndarray)): return [clean(v) for v in value]
    return str(value)


def response(operation: str, summary: dict[str, Any], *, artifacts: list[dict[str, Any]] | None = None, warnings: list[str] | None = None) -> dict[str, Any]:
    return clean({"success": True, "answer_ready": True, "operation": operation, "summary": summary, "evidence": [], "artifacts": artifacts or [], "warnings": warnings or [], "provenance": {"tool": operation, "version": "1.0.0"}, "recommended_next_tools": []})


def output_path(raw: str, *, directory: bool = False) -> Path:
    root = Path(os.environ.get("AI_DATASEEK_OUTPUT_ROOT", "/home/ubuntu/output")).resolve()
    path = Path(raw).resolve()
    if root != path and root not in path.parents: raise ValueError(f"output path must be below {root}")
    (path if directory else path.parent).mkdir(parents=True, exist_ok=True)
    return path


def artifact(path: Path, mime: str) -> dict[str, Any]:
    return {"path": str(path), "type": mime, "size_bytes": path.stat().st_size}


def select_slides(count: int, requested: list[int] | None, maximum: int) -> list[int]:
    if requested:
        selected=[]
        for slide in requested:
            if slide < 1 or slide > count: raise ValueError(f"slide is outside presentation: {slide}")
            if slide - 1 not in selected: selected.append(slide - 1)
        return selected[:maximum]
    if count <= maximum: return list(range(count))
    return sorted(set(np.linspace(0, count - 1, maximum, dtype=int).tolist()))


def notes_text(slide: Any) -> str:
    try:
        frame = slide.notes_slide.notes_text_frame
        return frame.text.strip() if frame is not None else ""
    except (AttributeError, ValueError):
        return ""


def slide_content(slide: Any, number: int) -> dict[str, Any]:
    texts=[]; tables=[]; charts=0; pictures=0
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip(): texts.append(shape.text.strip())
        if getattr(shape, "has_table", False):
            tables.append([[cell.text for cell in row.cells] for row in shape.table.rows])
        if getattr(shape, "has_chart", False): charts += 1
        if int(getattr(shape, "shape_type", 0)) == 13: pictures += 1
    title = slide.shapes.title.text.strip() if slide.shapes.title is not None and slide.shapes.title.text else None
    return {"slide": number, "title": title, "texts": texts, "notes": notes_text(slide), "tables": tables, "chart_count": charts, "picture_count": pictures}


def inspect_presentation(a: dict[str, Any]) -> dict[str, Any]:
    from pptx import Presentation
    path=Path(a["input_path"])
    if path.suffix.lower() not in {".pptx", ".pptm"}: raise ValueError("presentation inspection supports PPTX and PPTM")
    deck=Presentation(path); slides=[slide_content(slide, i + 1) for i, slide in enumerate(deck.slides)]
    summary={"name":path.name,"format":path.suffix.lower()[1:],"size_bytes":path.stat().st_size,"slide_count":len(slides),"width_inches":deck.slide_width/914400,"height_inches":deck.slide_height/914400,"slides":[{"slide":s["slide"],"title":s["title"],"text_characters":sum(len(t) for t in s["texts"]),"notes_characters":len(s["notes"]),"table_count":len(s["tables"]),"chart_count":s["chart_count"],"picture_count":s["picture_count"]} for s in slides]}
    return response("presentation_inspect", summary)


def extract_content(a: dict[str, Any]) -> dict[str, Any]:
    from pptx import Presentation
    deck=Presentation(a["input_path"]); indices=select_slides(len(deck.slides),a.get("slides"),100); limit=a.get("max_chars",200000); used=0; items=[]; warnings=[]
    for index in indices:
        item=slide_content(deck.slides[index],index+1); remaining=max(0,limit-used)
        for key in ("texts",):
            bounded=[]
            for text in item[key]:
                value=text[:remaining]; bounded.append(value); used+=len(value); remaining=max(0,limit-used)
                if remaining == 0: break
            item[key]=bounded
        item["notes"]=item["notes"][:remaining]; used+=len(item["notes"]); items.append(item)
        if used>=limit: warnings.append(f"presentation content truncated at {limit} characters"); break
    artifacts=[]
    if a.get("output_path"):
        path=output_path(a["output_path"])
        if path.suffix.lower()==".json": path.write_text(json.dumps(clean(items),ensure_ascii=False,indent=2),encoding="utf-8"); mime="application/json"
        elif path.suffix.lower() in {".md",".markdown"}:
            parts=[]
            for item in items:
                parts.append(f"## Slide {item['slide']}: {item['title'] or ''}".rstrip())
                parts.extend(item["texts"])
                if item["notes"]: parts.append(f"**Notes:** {item['notes']}")
                for table in item["tables"]: parts.append("```text\n"+"\n".join("\t".join(row) for row in table)+"\n```")
            path.write_text("\n\n".join(parts),encoding="utf-8"); mime="text/markdown"
        else: raise ValueError("presentation content output must be .json or .md")
        artifacts=[artifact(path,mime)]
    return response("pptx_extract_content",{"slides_extracted":len(items),"characters":used,"slides":items},artifacts=artifacts,warnings=warnings)


def render_slides(a: dict[str, Any]) -> dict[str, Any]:
    import pymupdf as fitz
    source=Path(a["input_path"]); directory=output_path(a["output_dir"],directory=True); artifacts=[]; rendered=[]
    with tempfile.TemporaryDirectory() as temp:
        completed=subprocess.run(["libreoffice","--headless","--convert-to","pdf","--outdir",temp,str(source)],capture_output=True,text=True,timeout=120)
        pdf=Path(temp)/(source.stem+".pdf")
        if completed.returncode or not pdf.exists(): raise ValueError("LibreOffice could not render the presentation")
        with fitz.open(pdf) as doc:
            indices=select_slides(len(doc),a.get("slides"),a.get("max_slides",6)); scale=a.get("dpi",144)/72
            for index in indices:
                pix=doc[index].get_pixmap(matrix=fitz.Matrix(scale,scale),alpha=False); path=directory/f"slide-{index+1}.png"; pix.save(path); artifacts.append(artifact(path,"image/png")); rendered.append({"slide":index+1,"width":pix.width,"height":pix.height,"path":str(path)})
    return response("pptx_render_slides",{"rendered":rendered},artifacts=artifacts)


FUNCTIONS={"presentation_inspect":inspect_presentation,"pptx_extract_content":extract_content,"pptx_render_slides":render_slides}


def main() -> int:
    parser=argparse.ArgumentParser(); parser.add_argument("tool"); parser.add_argument("payload"); args=parser.parse_args()
    try: result=FUNCTIONS[args.tool](json.loads(base64.urlsafe_b64decode(args.payload+"="*(-len(args.payload)%4))))
    except Exception as exc: result={"success":False,"answer_ready":True,"operation":args.tool,"error":f"{type(exc).__name__}: {exc}","warnings":[]}
    print(json.dumps(clean(result),ensure_ascii=False,allow_nan=False)); return 0 if result["success"] else 1


if __name__=="__main__": raise SystemExit(main())
